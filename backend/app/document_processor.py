"""Traitement des documents PDF, DOCX, PPTX pour le RAG"""

import os
import hashlib
import json
from pathlib import Path
from typing import List, Dict, Optional

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from .config import DOCUMENTS_DIR, CHUNK_SIZE, CHUNK_OVERLAP


class DocumentProcessor:
    """Charge les documents, extrait le texte et découpe en chunks."""

    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.txt', '.md'}

    def load_document(self, filepath: Path) -> Optional[str]:
        """Charge un document et retourne son texte brut."""
        ext = filepath.suffix.lower()
        
        if ext not in self.SUPPORTED_EXTENSIONS:
            print(f"Format non supporté: {ext} pour {filepath.name}")
            return None
        
        try:
            if ext == '.pdf':
                return self._load_pdf(filepath)
            elif ext == '.docx':
                return self._load_docx(filepath)
            elif ext == '.pptx':
                return self._load_pptx(filepath)
            elif ext in ('.txt', '.md'):
                return self._load_text(filepath)
        except Exception as e:
            print(f"Erreur lors du chargement de {filepath}: {e}")
            return None

    def _load_pdf(self, filepath: Path) -> str:
        """Extrait le texte d'un PDF."""
        if fitz is None:
            raise ImportError("PyMuPDF (fitz) n'est pas installé")
        text = ""
        doc = fitz.open(str(filepath))
        for page in doc:
            text += page.get_text() + "\n"
        doc.close()
        return text

    def _load_docx(self, filepath: Path) -> str:
        """Extrait le texte d'un DOCX."""
        if docx is None:
            raise ImportError("python-docx n'est pas installé")
        doc = docx.Document(str(filepath))
        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        # Ajouter les tableaux
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                if row_text.strip():
                    text += "\n" + row_text
        return text

    def _load_pptx(self, filepath: Path) -> str:
        """Extrait le texte d'un PPTX."""
        if Presentation is None:
            raise ImportError("python-pptx n'est pas installé")
        prs = Presentation(str(filepath))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text

    def _load_text(self, filepath: Path) -> str:
        """Charge un fichier texte."""
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()

    def chunk_text(self, text: str, metadata: Dict) -> List[Dict]:
        """Découpe le texte en chunks avec chevauchement."""
        chunks = []
        words = text.split()
        current_pos = 0
        
        while current_pos < len(words):
            chunk_words = words[current_pos:current_pos + CHUNK_SIZE]
            if not chunk_words:
                break
            
            chunk_text = " ".join(chunk_words)
            
            # Calculer un hash unique pour ce chunk
            chunk_hash = hashlib.md5(chunk_text.encode()).hexdigest()[:12]
            
            chunks.append({
                "id": f"{metadata.get('doc_id', 'doc')}_{chunk_hash}",
                "text": chunk_text,
                "metadata": {**metadata, "chunk_index": current_pos // CHUNK_SIZE}
            })
            
            current_pos += CHUNK_SIZE - CHUNK_OVERLAP
        
        return chunks

    def process_all_documents(self) -> List[Dict]:
        """Traite tous les documents du dossier data/documents."""
        all_chunks = []
        
        if not DOCUMENTS_DIR.exists():
            print(f"Le dossier {DOCUMENTS_DIR} n'existe pas")
            return all_chunks
        
        # Parcourir les sous-dossiers (classe > chapitre)
        for classe_dir in sorted(DOCUMENTS_DIR.iterdir()):
            if not classe_dir.is_dir():
                continue
            
            classe_name = classe_dir.name
            
            for chapitre_dir in sorted(classe_dir.iterdir()):
                if not chapitre_dir.is_dir():
                    # Aussi traiter les fichiers directement dans le dossier classe
                    if chapitre_dir.is_file():
                        text = self.load_document(chapitre_dir)
                        if text and text.strip():
                            doc_id = f"{classe_name}_{chapitre_dir.stem}"
                            metadata = {
                                "doc_id": doc_id,
                                "classe": classe_name,
                                "chapitre": "Général",
                                "filename": chapitre_dir.name,
                                "source": str(chapitre_dir)
                            }
                            chunks = self.chunk_text(text, metadata)
                            if chunks:
                                all_chunks.extend(chunks)
                                print(f"[OK] {chapitre_dir.name}: {len(chunks)} chunks générés")
                    continue
                
                chapitre_name = chapitre_dir.name
                
                for filepath in chapitre_dir.iterdir():
                    if not filepath.is_file():
                        continue
                    
                    text = self.load_document(filepath)
                    if text is None or not text.strip():
                        continue
                    
                    doc_id = f"{classe_name}_{chapitre_name}_{filepath.stem}"
                    metadata = {
                        "doc_id": doc_id,
                        "classe": classe_name,
                        "chapitre": chapitre_name,
                        "filename": filepath.name,
                        "source": str(filepath)
                    }
                    
                    chunks = self.chunk_text(text, metadata)
                    all_chunks.extend(chunks)
                    
                    print(f"[OK] {filepath.name}: {len(chunks)} chunks générés")
        
        return all_chunks

    def get_document_stats(self) -> Dict:
        """Retourne des statistiques sur les documents chargés."""
        stats = {
            "total_files": 0,
            "by_classe": {},
            "by_format": {ext: 0 for ext in self.SUPPORTED_EXTENSIONS}
        }
        
        if not DOCUMENTS_DIR.exists():
            return stats
        
        for item in sorted(DOCUMENTS_DIR.iterdir()):
            if not item.is_dir():
                continue
            classe_name = item.name
            stats["by_classe"][classe_name] = {"files": 0, "chapitres": 0}
            
            for sub_item in sorted(item.iterdir()):
                if not sub_item.is_dir():
                    # Count files directly in class dir
                    if sub_item.is_file():
                        stats["total_files"] += 1
                        stats["by_classe"][classe_name]["files"] += 1
                        ext = sub_item.suffix.lower()
                        if ext in stats["by_format"]:
                            stats["by_format"][ext] += 1
                    continue
                stats["by_classe"][classe_name]["chapitres"] += 1
                
                for filepath in sub_item.iterdir():
                    if filepath.is_file():
                        stats["total_files"] += 1
                        stats["by_classe"][classe_name]["files"] += 1
                        ext = filepath.suffix.lower()
                        if ext in stats["by_format"]:
                            stats["by_format"][ext] += 1
        
        return stats
